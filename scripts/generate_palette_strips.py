#!/usr/bin/env python3
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

PALETTES = {
    'nature_blue': ['F7FAFC','1F4E79','2C7FB8','7FCDBB','EDF8FB','263238'],
    'chinese_science_blue': ['F8FAFC','163B6D','2F6F9F','7BAFD4','EAF3FA','202833'],
    'sci_cjk_bilingual': ['FAFAF7','1E3A5F','357ABD','B55A30','EFE8D8','222222'],
    'science_graphite': ['F7F7F5','2B2D30','5C6670','9AA4AE','ECECE8','242424'],
    'advanced_materials_teal': ['F6FBFA','0E5A5A','1C8C8C','8FD0C9','E7F4F1','1F2A2A'],
    'engineering_amber': ['FBFAF5','4B3B1F','B97913','E2B84F','F5ECD0','241C12'],
    'ai4science_indigo': ['F8F8FD','232A7A','4C63D2','8EA0F2','ECEEFF','1D2030'],
    'minimal_mono': ['FFFFFF','222222','666666','AAAAAA','F3F3F3','111111'],
}

def rgb(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def text(slide,x,y,w,h,s,size=10,bold=False):
    box = slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=box.text_frame; tf.clear(); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
    r=p.add_run(); r.text=s; r.font.size=Pt(size); r.font.bold=bold; r.font.name='Arial'; r.font.color.rgb=rgb('222222')

def main():
    ap=argparse.ArgumentParser(); ap.add_argument('out'); args=ap.parse_args()
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    text(slide,0.5,0.25,12,0.35,'Palette strip preview',20,True)
    y=0.85
    for name, colors in PALETTES.items():
        text(slide,0.55,y+0.08,2.15,0.25,name,9.5,True)
        x=2.75
        for c in colors:
            rect=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(1.25),Inches(0.32))
            rect.fill.solid(); rect.fill.fore_color.rgb=rgb(c); rect.line.color.rgb=rgb('DDDDDD')
            text(slide,x,y+0.36,1.2,0.18,'#'+c,6.5,False)
            x+=1.3
        y+=0.72
    prs.save(args.out)
    print(args.out)
if __name__ == '__main__':
    main()
