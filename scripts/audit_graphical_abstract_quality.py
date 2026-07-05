#!/usr/bin/env python3
"""Publication-grade quality audit for graphical abstract PPTX outputs."""
import json, sys, zipfile, re, math
from pathlib import Path

try:
    from pptx import Presentation
except Exception:
    Presentation = None

APPROVED_PALETTES = {"nature_blue", "science_graphite", "advanced_materials_teal", "cell_biomedical", "engineering_amber", "ai4science_indigo", "earth_environment", "minimal_mono", "chinese_science_blue", "chinese_academy_red", "sci_cjk_bilingual"}
CONTENT_FIELDS=['main_process','visual_objects','method_or_model','mechanism','key_result','application','source_figures_to_redraw','forbidden_content']


def cjk_count(s):
    return len(re.findall(r"[\u4e00-\u9fff]", s or ""))

def text_len(v):
    if isinstance(v,str): return len(v.strip())
    if isinstance(v,list): return sum(text_len(x) for x in v)
    if isinstance(v,dict): return sum(text_len(x) for x in v.values())
    return 0

def content_score(ac):
    if not ac: return 0
    if isinstance(ac,str): return 2 if len(ac.strip()) >= 100 else (1 if len(ac.strip()) >= 20 else 0)
    score=0
    for f in ['main_process','visual_objects','method_or_model','key_result','application']:
        if text_len(ac.get(f))>0: score += 1
    if text_len(ac.get('mechanism'))>0: score += 1
    if text_len(ac.get('source_figures_to_redraw'))>0: score += 1
    if text_len(ac.get('forbidden_content'))>0: score += 1
    return min(score,8)

def pptx_zip_issues(pptx):
    issues=[]
    with zipfile.ZipFile(pptx,'r') as zf:
        for name in zf.namelist():
            low=name.lower()
            if low == 'docprops/thumbnail.jpeg':
                continue
            if low.startswith('ppt/media/') or low.startswith('ppt/embeddings/'):
                issues.append(f'Forbidden embedded asset: {name}')
            if Path(low).suffix in {'.png','.jpg','.jpeg','.gif','.bmp','.tif','.tiff','.svg','.wmf','.emf','.mp4','.mov','.avi','.wav','.mp3'}:
                issues.append(f'Forbidden embedded asset extension: {name}')
            if low.endswith('.xml') or low.endswith('.rels'):
                data=zf.read(name)
                for pat,label in [(rb'<p:pic\b','picture object'),(rb'<pic:pic\b','picture object'),(rb'oleObject','OLE object'),(rb'TargetMode="External"','external relationship'),(rb'videoFile','video'),(rb'audioFile','audio')]:
                    if re.search(pat,data): issues.append(f'Detected {label} in {name}')
    return sorted(set(issues))

def overlap_area(a,b):
    l=max(a[0],b[0]); t=max(a[1],b[1]); r=min(a[2],b[2]); bot=min(a[3],b[3])
    return max(0,r-l)*max(0,bot-t)

def main():
    if len(sys.argv) not in (3,4):
        print('Usage: audit_graphical_abstract_quality.py <spec.json> <file.pptx> [report.json]')
        return 2
    spec_path=Path(sys.argv[1]); pptx=Path(sys.argv[2])
    spec=json.loads(spec_path.read_text(encoding='utf-8'))
    issues=[]; warnings=[]; metrics={}
    strict=bool(spec.get('strict_mode') or spec.get('journal_profile') in {'international_top_journal','chinese_top_journal'})

    sc=content_score(spec.get('abstract_content'))
    metrics['content_sufficiency_score']=sc
    if sc==0: issues.append('Q0 failed: missing abstract_content')
    elif sc<5:
        msg=f'Content brief is weak for strict generation (score={sc}/8)'
        (issues if strict else warnings).append(msg)

    if spec.get('palette_name') not in APPROVED_PALETTES:
        issues.append('Q5 failed: palette_name is not approved')
    if not spec.get('central_claim') and not spec.get('visual_claim'):
        warnings.append('Central visual claim is absent')
    if len(spec.get('blocks',[]) or []) > 7:
        warnings.append('More than 7 blocks may be too dense')
    if len(spec.get('blocks',[]) or []) < 3:
        warnings.append('Fewer than 3 blocks may be under-developed')
    for b in spec.get('blocks',[]) or []:
        if not b.get('claim_source'):
            (issues if strict else warnings).append(f"Missing claim_source for block {b.get('id')}")
        if spec.get('language') in {'zh-CN','zh-TW','bilingual'} and cjk_count(str(b.get('label',''))) > 14:
            warnings.append(f"Chinese block label may be too long: {b.get('label')}")
        if spec.get('language') == 'en' and len(str(b.get('label',''))) > 46:
            warnings.append(f"English block label may be too long: {b.get('label')}")
        if text_len(b.get('bullets',[])) > 140:
            warnings.append(f"Dense bullet text in block {b.get('id')}")

    try:
        issues.extend(pptx_zip_issues(pptx))
    except Exception as e:
        issues.append(f'Cannot inspect PPTX zip: {e}')

    if Presentation is not None and pptx.exists():
        try:
            prs=Presentation(str(pptx))
            shape_count=0; text_shapes=0; long_text=0; out_bounds=0; overlaps=0; text_boxes=[]
            sw=prs.slide_width; sh=prs.slide_height
            for slide_idx, slide in enumerate(prs.slides):
                boxes=[]
                for shape in slide.shapes:
                    shape_count += 1
                    left,top,right,bottom=shape.left,shape.top,shape.left+shape.width,shape.top+shape.height
                    if left < 0 or top < 0 or right > sw or bottom > sh:
                        out_bounds += 1
                    if getattr(shape,'has_text_frame',False) and shape.has_text_frame:
                        txt=' '.join(p.text for p in shape.text_frame.paragraphs)
                        if txt.strip():
                            text_shapes += 1
                            if slide_idx == 0:
                                boxes.append((left,top,right,bottom,txt))
                                if len(txt) > 180: long_text += 1
                if slide_idx == 0:
                    for i in range(len(boxes)):
                        for j in range(i+1,len(boxes)):
                            a=boxes[i]; b=boxes[j]
                            area=overlap_area(a,b)
                            if area and area > 0.85*min((a[2]-a[0])*(a[3]-a[1]),(b[2]-b[0])*(b[3]-b[1])):
                                overlaps += 1
            metrics.update({'shape_count':shape_count,'text_shapes':text_shapes,'long_text_shapes':long_text,'out_of_bounds_shapes':out_bounds,'text_overlap_pairs':overlaps,'slide_count':len(prs.slides)})
            if out_bounds: issues.append(f'Q2 failed: {out_bounds} shapes exceed slide bounds')
            if overlaps > 10: warnings.append(f'{overlaps} text-box overlap pairs detected; inspect layout manually')
            if long_text: warnings.append(f'{long_text} text shapes are long; graphical abstracts should avoid paragraphs')
            if shape_count < 15: warnings.append('Low shape count; output may be under-developed')
            if text_shapes < 4: warnings.append('Few editable text elements detected')
        except Exception as e:
            warnings.append(f'PPTX semantic inspection unavailable: {e}')

    score=max(0,100-12*len(issues)-4*len(warnings))
    report={'score':score,'strict_mode':strict,'issues':issues,'warnings':warnings,'metrics':metrics}
    print(json.dumps(report,ensure_ascii=False,indent=2))
    if len(sys.argv)==4:
        Path(sys.argv[3]).write_text(json.dumps(report,ensure_ascii=False,indent=2),encoding='utf-8')
    return 1 if issues else 0

if __name__=='__main__':
    raise SystemExit(main())
