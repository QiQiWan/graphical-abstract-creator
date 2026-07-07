#!/usr/bin/env python3
import argparse, json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

PALETTES = {
    'nature_blue': ['F7FAFC','1F4E79','2C7FB8','7FCDBB','EDF8FB','263238'],
    'chinese_science_blue': ['F8FAFC','163B6D','2F6F9F','7BAFD4','EAF3FA','202833'],
    'sci_cjk_bilingual': ['FAFAF7','1E3A5F','357ABD','B55A30','EFE8D8','222222'],
    'science_graphite': ['F7F7F5','2B2D30','5C6670','9AA4AE','ECECE8','242424'],
    'advanced_materials_teal': ['F6FBFA','0E5A5A','1C8C8C','8FD0C9','E7F4F1','1F2A2A'],
    'engineering_amber': ['FBFAF5','4B3B1F','B97913','E2B84F','F5ECD0','241C12'],
    'ai4science_indigo': ['F8F8FD','232A7A','4C63D2','8EA0F2','ECEEFF','1D2030'],
    'minimal_mono': ['FFFFFF','222222','666666','AAAAAA','F3F3F3','111111'],
}

def rgb(hexs):
    hexs = hexs.strip().lstrip('#')
    return RGBColor(int(hexs[0:2],16), int(hexs[2:4],16), int(hexs[4:6],16))

def add_textbox(slide, x,y,w,h,text, size=12, bold=False, color='222222', align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = 'Arial'
    r.font.color.rgb = rgb(color)
    return box

def add_module(slide, block, x,y,w,h, pal, core=False):
    bg, primary, secondary, accent, light, text = pal
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(light if not core else 'FFFFFF')
    shape.line.color.rgb = rgb(primary if core else secondary)
    shape.line.width = Pt(1.6 if core else 0.9)
    add_textbox(slide, x+0.12, y+0.12, w-0.24, 0.28, block.get('title',''), size=15 if core else 12.5, bold=True, color=primary, align=PP_ALIGN.CENTER)
    labels = block.get('labels', [])[:2]
    for i, lab in enumerate(labels):
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.18), Inches(y+0.55+i*0.38), Inches(w-0.36), Inches(0.26))
        pill.fill.solid(); pill.fill.fore_color.rgb = rgb('FFFFFF')
        pill.line.color.rgb = rgb('D7DEE8')
        add_textbox(slide, x+0.24, y+0.585+i*0.38, w-0.48, 0.16, str(lab), size=9.2, color=text, align=PP_ALIGN.CENTER)
    # simple vector glyph
    cx = x + w/2
    gy = y + h - 0.62
    if block.get('role') == 'method' or core:
        for j in range(3):
            node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx-0.45+j*0.45), Inches(gy), Inches(0.16), Inches(0.16))
            node.fill.solid(); node.fill.fore_color.rgb = rgb(secondary)
            node.line.color.rgb = rgb(secondary)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cx-0.37), Inches(gy+0.08), Inches(cx+0.53), Inches(gy+0.08))
        line.line.color.rgb = rgb(secondary); line.line.width = Pt(1.0)
    elif block.get('role') == 'input':
        for j in range(3):
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx-0.42+j*0.3), Inches(gy), Inches(0.22), Inches(0.16))
            rect.fill.solid(); rect.fill.fore_color.rgb = rgb('FFFFFF')
            rect.line.color.rgb = rgb(secondary)
    else:
        for j in range(3):
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx-0.38+j*0.28), Inches(gy+0.18-j*0.06), Inches(0.16), Inches(0.16+j*0.06))
            bar.fill.solid(); bar.fill.fore_color.rgb = rgb(accent)
            bar.line.color.rgb = rgb(accent)
    return shape

def positions(blocks):
    n = len(blocks)
    if n <= 3:
        return [(0.9,2.3,2.35,2.0),(4.05,1.85,3.0,2.65),(8.1,2.3,2.35,2.0)][:n]
    if n == 4:
        return [(0.65,2.35,2.25,1.9),(3.35,1.92,2.75,2.55),(6.55,2.35,2.25,1.9),(9.25,2.35,2.25,1.9)]
    return [(0.45,2.35,2.05,1.85),(2.85,2.35,2.05,1.85),(5.15,1.9,2.7,2.55),(8.15,2.35,2.05,1.85),(10.55,2.35,2.05,1.85)]

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('spec')
    ap.add_argument('out')
    args = ap.parse_args()
    data = json.loads(Path(args.spec).read_text(encoding='utf-8'))
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pal = PALETTES.get(data.get('palette_name','nature_blue'), PALETTES['nature_blue'])
    bg, primary, secondary, accent, light, text = pal
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(bg)
    add_textbox(slide, 0.55,0.28,12.1,0.46, data.get('title','Graphical abstract'), size=23, bold=True, color=primary, align=PP_ALIGN.CENTER)
    subtitle = data.get('subtitle','')
    if subtitle:
        add_textbox(slide,0.8,0.78,11.7,0.28,subtitle,size=10.5,color=text,align=PP_ALIGN.CENTER)
    blocks = data.get('blocks', [])
    pos = positions(blocks)
    centers=[]
    for block, p in zip(blocks, pos):
        core = block.get('emphasis') == 'core' or block.get('role') == 'method'
        add_module(slide, block, *p, pal, core=core)
        centers.append((p[0]+p[2], p[1]+p[3]/2, p[0], p[1]+p[3]/2))
    for i in range(len(centers)-1):
        x1,y1,_,_ = centers[i]
        _,_,x2,y2 = centers[i+1]
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1+0.18), Inches(y1), Inches(x2-0.18), Inches(y2))
        conn.line.color.rgb = rgb(primary); conn.line.width = Pt(1.7)
        conn.line.end_arrowhead = True
    msg = data.get('key_message','')
    if msg:
        strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(6.35), Inches(11.25), Inches(0.48))
        strip.fill.solid(); strip.fill.fore_color.rgb = rgb('FFFFFF')
        strip.line.color.rgb = rgb('D7DEE8')
        add_textbox(slide,1.15,6.47,11.05,0.18,msg,size=11.5,color=text,align=PP_ALIGN.CENTER)
    prs.save(args.out)
    print(args.out)
if __name__ == '__main__':
    main()
